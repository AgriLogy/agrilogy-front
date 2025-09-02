from pptx import Presentation

# Create new presentation
prs = Presentation()

# Define structured slides
slides = [
    ("Smart Agriculture with AI", "Predicting Environment and Diagnosing Plant Health\nZakaria Makhkhas - 2025"),
    ("Problem Statement", "• Unpredictable weather affects crop yield\n• Farmers lack early disease detection tools\n• No system to predict fruit growth\n\n✅ Need: Smart system using sensors and AI"),
    ("Project Objectives", "1. Predict weather, humidity, and fruit size using historical sensor data\n2. Classify leaf health based on image analysis\n3. Create an interactive monitoring dashboard"),
    ("System Architecture", "• Two main components:\n  - Time-Series Forecasting for environment and growth\n  - Image Classification for plant health\n• Integrated with frontend dashboard and backend API"),
    ("Time-Series Forecasting", "📍 Inputs: Temperature, Humidity, NPK, Fruit size (historical)\n🧠 Model: LSTM / Random Forest\n🎯 Output: Next values prediction\n💡 Use: Alerting and planning irrigation/fertilization"),
    ("Leaf Image Classification", "📍 Inputs: Images from smartphone or camera\n🧠 Model: CNN (ResNet/MobileNet)\n🎯 Output: Healthy / Unhealthy / Disease category\n📸 Use: Early detection and treatment"),
    ("Datasets", "• Captor Data: Collected from farm sensors over time\n• Leaf Images: PlantVillage dataset or own camera captures\n• Storage: Structured in database for ML pipeline"),
    ("Technologies", "• ML/DL: Scikit-learn, TensorFlow, PyTorch\n• Backend: Flask or FastAPI\n• Frontend: React + Chakra UI\n• Deployment: Docker, Render, or local VPS"),
    ("Expected Outcomes", "• Real-time environmental forecasts\n• Visual disease diagnosis\n• Centralized dashboard for farmers\n• Improved yield and reduced crop loss"),
    ("Future Work", "• Real-time alerts via SMS/email\n• Expand models to cover more crops/diseases\n• Add soil moisture and irrigation control module"),
    ("Conclusion", "AI-powered farming can transform agriculture.\nThis project helps farmers make better decisions.\nWe seek support to develop, test, and deploy the system.")
]

# Add slides to presentation
for title, content in slides:
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Save file
pptx_file_path = "/mnt/data/Smart_Agriculture_AI_Redesigned.pptx"
prs.save(pptx_file_path)
pptx_file_path
